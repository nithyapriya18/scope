#!/usr/bin/env python3
"""
Physical Document Generator Script
Generates Word, Excel, and PowerPoint files from structured JSON content
"""

import sys
import json
import argparse
from pathlib import Path
from datetime import datetime
import psycopg2
from psycopg2.extras import RealDictCursor

# Document generation libraries
try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
except ImportError:
    print("ERROR: python-docx not installed. Run: pip install python-docx")
    sys.exit(1)

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
except ImportError:
    print("ERROR: openpyxl not installed. Run: pip install openpyxl")
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches as PptInches, Pt as PptPt
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


# Database configuration
DB_CONFIG = {
    'dbname': 'verity_scope',
    'user': 'nithya',
    'password': 'dev123',
    'host': '127.0.0.1',
    'port': 5432
}

# Output directory
OUTPUT_DIR = Path(__file__).parent.parent / 'generated_documents'
OUTPUT_DIR.mkdir(exist_ok=True)

# PetaSight brand colors
BRAND_PRIMARY = RGBColor(218, 54, 92)  # #da365c magenta
BRAND_SECONDARY = RGBColor(30, 64, 175)  # #1e40af blue
BRAND_ACCENT = RGBColor(52, 211, 153)  # #34d399 emerald


def get_db_connection():
    """Create database connection"""
    return psycopg2.connect(**DB_CONFIG, cursor_factory=RealDictCursor)


def fetch_document_data(document_id):
    """Fetch document content from database"""
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT
                    d.id,
                    d.document_type,
                    d.document_content,
                    d.pricing_pack_id,
                    p.total_price,
                    o.rfp_title,
                    o.client_name
                FROM documents d
                JOIN pricing_packs p ON d.pricing_pack_id = p.id
                JOIN wbs w ON p.wbs_id = w.id
                JOIN hcp_shortlists h ON w.hcp_shortlist_id = h.id
                JOIN sample_plans sp ON h.sample_plan_id = sp.id
                JOIN scopes s ON sp.scope_id = s.id
                JOIN briefs b ON s.brief_id = b.id
                JOIN opportunities o ON b.opportunity_id = o.id
                WHERE d.id = %s
            """, (document_id,))
            return cur.fetchone()
    finally:
        conn.close()


def generate_proposal_docx(content, output_path, metadata):
    """Generate Proposal Word document"""
    print(f"📄 Generating Proposal DOCX: {output_path}")

    doc = Document()

    # Set document margins
    sections = doc.sections
    for section in sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)

    # Cover Page
    title = doc.add_heading('Proposal', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.runs[0].font.color.rgb = BRAND_PRIMARY

    cover = content.get('proposalContent', {}).get('coverPage', {})
    doc.add_paragraph(f"\n{cover.get('rfpTitle', metadata['rfp_title'])}", style='Title')
    doc.add_paragraph(f"\nPrepared for:", style='Heading 2')
    doc.add_paragraph(f"{cover.get('clientName', metadata['client_name'])}")
    doc.add_paragraph(f"\nPrepared by:", style='Heading 2')
    doc.add_paragraph("PetaSight - Verity Scope")
    doc.add_paragraph("Pharmaceutical Market Research Excellence")
    doc.add_paragraph(f"\n{cover.get('date', datetime.now().strftime('%B %d, %Y'))}")
    doc.add_paragraph("\n\nCONFIDENTIAL", style='Heading 2').alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_page_break()

    # Executive Summary
    doc.add_heading('Executive Summary', 1)
    exec_summary = content.get('proposalContent', {}).get('executiveSummary', {})
    doc.add_paragraph(exec_summary.get('understanding', 'We understand your research needs...'))
    doc.add_paragraph()
    doc.add_heading('Our Approach', 2)
    doc.add_paragraph(exec_summary.get('approach', 'We recommend a qualitative approach...'))
    doc.add_paragraph()
    doc.add_heading('Why PetaSight?', 2)
    doc.add_paragraph(exec_summary.get('value', 'Our value propositions...'))
    doc.add_page_break()

    # Company Background
    doc.add_heading('Company Background', 1)
    company = content.get('proposalContent', {}).get('companyBackground', {})
    doc.add_paragraph(company.get('overview', 'PetaSight is a leading pharmaceutical market research firm...'))
    doc.add_page_break()

    # Understanding of Requirements
    doc.add_heading('Understanding of Requirements', 1)
    understanding = content.get('proposalContent', {}).get('understanding', {})
    doc.add_paragraph(understanding.get('objectives', 'Research objectives as we understand them...'))
    doc.add_page_break()

    # Proposed Methodology
    doc.add_heading('Proposed Methodology', 1)
    methodology = content.get('proposalContent', {}).get('methodology', {})
    doc.add_paragraph(methodology.get('studyDesign', 'Study design details...'))
    doc.add_page_break()

    # Sample Plan
    doc.add_heading('Sample Plan', 1)
    sample_plan = content.get('proposalContent', {}).get('samplePlan', {})
    doc.add_paragraph(sample_plan.get('size', 'Sample size and distribution...'))
    doc.add_page_break()

    # Project Team
    doc.add_heading('Project Team', 1)
    team = content.get('proposalContent', {}).get('team', {})
    doc.add_paragraph(team.get('pm', {}).get('bio', 'Project Manager biography...'))
    doc.add_page_break()

    # Timeline
    doc.add_heading('Timeline', 1)
    timeline = content.get('proposalContent', {}).get('timeline', {})
    doc.add_paragraph(str(timeline.get('phases', 'Project phases and milestones...')))
    doc.add_page_break()

    # Pricing Summary
    doc.add_heading('Pricing Summary', 1)
    pricing = content.get('proposalContent', {}).get('pricing', {})
    doc.add_paragraph(f"Total Investment: {pricing.get('total', '$0')}")
    doc.add_paragraph(f"Payment Terms: {pricing.get('terms', 'TBD')}")

    # Save
    doc.save(output_path)
    print(f"✅ Proposal saved: {output_path}")


def generate_sow_docx(content, output_path, metadata):
    """Generate Statement of Work Word document"""
    print(f"📄 Generating SOW DOCX: {output_path}")

    doc = Document()

    # Title
    title = doc.add_heading('Statement of Work', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.runs[0].font.color.rgb = BRAND_PRIMARY

    # Project Overview
    doc.add_heading('Project Overview', 1)
    sow = content.get('sowContent', {})
    overview = sow.get('projectOverview', {})
    doc.add_paragraph(f"Project Title: {overview.get('title', metadata['rfp_title'])}")
    doc.add_paragraph(f"Client: {overview.get('client', metadata['client_name'])}")
    doc.add_paragraph(f"Start Date: {overview.get('startDate', 'TBD')}")
    doc.add_paragraph(f"End Date: {overview.get('endDate', 'TBD')}")
    doc.add_page_break()

    # Scope of Services
    doc.add_heading('Scope of Services', 1)
    scope = sow.get('scopeOfServices', {})
    doc.add_paragraph(str(scope.get('tasks', 'Detailed task list...')))
    doc.add_page_break()

    # Deliverables
    doc.add_heading('Deliverables', 1)
    deliverables = sow.get('deliverables', {})
    doc.add_paragraph(str(deliverables.get('list', 'Itemized deliverables...')))
    doc.add_page_break()

    # Client Responsibilities
    doc.add_heading('Client Responsibilities', 1)
    client_resp = sow.get('clientResponsibilities', {})
    doc.add_paragraph(str(client_resp.get('items', 'Client responsibilities...')))
    doc.add_page_break()

    # Terms and Conditions
    doc.add_heading('Terms and Conditions', 1)
    terms = sow.get('terms', {})
    doc.add_paragraph(terms.get('payment', 'Payment terms: 50/25/25'))
    doc.add_paragraph(terms.get('confidentiality', 'Confidentiality clause...'))

    # Save
    doc.save(output_path)
    print(f"✅ SOW saved: {output_path}")


def generate_pricing_xlsx(content, output_path, metadata):
    """Generate Pricing Pack Excel workbook"""
    print(f"📊 Generating Pricing Pack XLSX: {output_path}")

    wb = Workbook()

    # Remove default sheet
    wb.remove(wb.active)

    # Tab 1: Summary
    ws1 = wb.create_sheet("Summary")
    pricing_content = content.get('pricingPackContent', {})
    summary = pricing_content.get('summary', {})

    ws1['A1'] = 'Pricing Summary'
    ws1['A1'].font = Font(size=16, bold=True, color='DA365C')
    ws1['A3'] = 'Total Project Price:'
    ws1['B3'] = summary.get('totalPrice', '$0')
    ws1['A4'] = 'Payment Schedule:'
    ws1['B4'] = summary.get('paymentSchedule', '50/25/25')
    ws1['A5'] = 'RFP Budget:'
    ws1['B5'] = summary.get('rfpBudget', 'TBD')

    # Tab 2: Detailed Breakdown
    ws2 = wb.create_sheet("Detailed Breakdown")
    breakdown = pricing_content.get('breakdown', {})

    ws2['A1'] = 'Detailed Cost Breakdown'
    ws2['A1'].font = Font(size=16, bold=True, color='DA365C')

    headers = ['Category', 'Description', 'Hours/Units', 'Rate', 'Amount']
    for col, header in enumerate(headers, 1):
        cell = ws2.cell(3, col)
        cell.value = header
        cell.font = Font(bold=True)
        cell.fill = PatternFill(start_color='DA365C', end_color='DA365C', fill_type='solid')

    row = 4
    for item in breakdown.get('items', []):
        ws2.cell(row, 1).value = item.get('category', '')
        ws2.cell(row, 2).value = item.get('description', '')
        ws2.cell(row, 3).value = item.get('units', 0)
        ws2.cell(row, 4).value = item.get('rate', '$0')
        ws2.cell(row, 5).value = item.get('amount', '$0')
        row += 1

    # Tab 3: Assumptions
    ws3 = wb.create_sheet("Assumptions")
    assumptions = pricing_content.get('assumptions', {})

    ws3['A1'] = 'Pricing Assumptions'
    ws3['A1'].font = Font(size=16, bold=True, color='DA365C')
    ws3['A3'] = 'Inclusions:'
    ws3['A4'] = str(assumptions.get('inclusions', 'TBD'))
    ws3['A6'] = 'Exclusions:'
    ws3['A7'] = str(assumptions.get('exclusions', 'TBD'))
    ws3['A9'] = 'Notes:'
    ws3['A10'] = str(assumptions.get('notes', 'TBD'))

    # Save
    wb.save(output_path)
    print(f"✅ Pricing Pack saved: {output_path}")


def generate_presentation_pptx(content, output_path, metadata):
    """Generate Capabilities Presentation PowerPoint"""
    print(f"📊 Generating Presentation PPTX: {output_path}")

    prs = Presentation()
    prs.slide_width = PptInches(10)
    prs.slide_height = PptInches(7.5)

    pres_content = content.get('presentationContent', {})

    # Slide 1: Cover
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    title_box = slide1.shapes.add_textbox(PptInches(1), PptInches(2.5), PptInches(8), PptInches(1))
    title_frame = title_box.text_frame
    title_frame.text = f"Proposal for {metadata['rfp_title']}"
    title_frame.paragraphs[0].font.size = PptPt(44)
    title_frame.paragraphs[0].font.bold = True

    subtitle_box = slide1.shapes.add_textbox(PptInches(1), PptInches(4), PptInches(8), PptInches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "PetaSight - Verity Scope"
    subtitle_frame.paragraphs[0].font.size = PptPt(24)

    # Slide 2: Agenda
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    slide2.shapes.title.text = "Agenda"
    content_box = slide2.placeholders[1]
    tf = content_box.text_frame
    tf.text = "• Understanding Your Needs\n• Our Approach\n• Sample Design\n• Timeline\n• Our Team\n• Why PetaSight?\n• Pricing\n• Next Steps"

    # Slide 3: Understanding Your Needs
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Understanding Your Needs"
    understanding = pres_content.get('understanding', {})
    content_box = slide3.placeholders[1]
    content_box.text_frame.text = understanding.get('objectives', 'Research objectives...')

    # Slide 4: Our Approach
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Our Approach"
    approach = pres_content.get('approach', {})
    content_box = slide4.placeholders[1]
    content_box.text_frame.text = approach.get('methodology', 'Recommended methodology...')

    # Slide 5: Pricing Summary
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Investment"
    pricing = pres_content.get('pricing', {})
    content_box = slide5.placeholders[1]
    content_box.text_frame.text = f"Total Investment: {pricing.get('total', '$0')}\n\nPayment Terms: {pricing.get('terms', '50/25/25')}"

    # Slide 6: Next Steps
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Next Steps"
    next_steps = pres_content.get('nextSteps', {})
    content_box = slide6.placeholders[1]
    content_box.text_frame.text = "1. Review proposal\n2. Approval\n3. Kickoff meeting\n4. Project execution"

    # Save
    prs.save(output_path)
    print(f"✅ Presentation saved: {output_path}")


def update_document_file_path(document_id, file_path):
    """Update document record with generated file path"""
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                UPDATE documents
                SET file_path = %s, updated_at = now()
                WHERE id = %s
            """, (str(file_path), document_id))
        conn.commit()
    finally:
        conn.close()


def generate_document(document_id):
    """Generate physical document from database content"""
    print(f"\n🚀 Generating document ID: {document_id}")

    # Fetch document data
    doc_data = fetch_document_data(document_id)
    if not doc_data:
        print(f"❌ Document not found: {document_id}")
        return False

    doc_type = doc_data['document_type']
    content = doc_data['document_content']

    metadata = {
        'rfp_title': doc_data['rfp_title'],
        'client_name': doc_data['client_name']
    }

    # Generate filename
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    safe_client = metadata['client_name'].replace(' ', '_') if metadata['client_name'] else 'Client'

    try:
        if doc_type == 'proposal':
            filename = f"{safe_client}_Proposal_{timestamp}.docx"
            output_path = OUTPUT_DIR / filename
            generate_proposal_docx(content, output_path, metadata)

        elif doc_type == 'sow':
            filename = f"{safe_client}_SOW_{timestamp}.docx"
            output_path = OUTPUT_DIR / filename
            generate_sow_docx(content, output_path, metadata)

        elif doc_type == 'pricing_pack':
            filename = f"{safe_client}_Pricing_{timestamp}.xlsx"
            output_path = OUTPUT_DIR / filename
            generate_pricing_xlsx(content, output_path, metadata)

        elif doc_type == 'presentation':
            filename = f"{safe_client}_Presentation_{timestamp}.pptx"
            output_path = OUTPUT_DIR / filename
            generate_presentation_pptx(content, output_path, metadata)

        else:
            print(f"❌ Unknown document type: {doc_type}")
            return False

        # Update database with file path
        update_document_file_path(document_id, output_path)

        print(f"\n✅ SUCCESS: Document generated at {output_path}\n")
        return True

    except Exception as e:
        print(f"❌ ERROR generating {doc_type}: {e}")
        import traceback
        traceback.print_exc()
        return False


def main():
    parser = argparse.ArgumentParser(description='Generate physical documents from Verity Scope database')
    parser.add_argument('document_ids', nargs='+', help='Document IDs to generate')
    args = parser.parse_args()

    print("=" * 80)
    print("📄 Verity Scope Document Generator")
    print("=" * 80)

    success_count = 0
    for doc_id in args.document_ids:
        if generate_document(doc_id):
            success_count += 1

    print("=" * 80)
    print(f"✅ Generated {success_count}/{len(args.document_ids)} documents")
    print(f"📁 Output directory: {OUTPUT_DIR}")
    print("=" * 80)

    return 0 if success_count == len(args.document_ids) else 1


if __name__ == '__main__':
    sys.exit(main())
